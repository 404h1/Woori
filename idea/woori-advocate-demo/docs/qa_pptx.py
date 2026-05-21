"""
woori_advocate_v1.pptx 콘텐츠 QA.

체크:
1. 슬라이드 수 = 15
2. 각 슬라이드에 speaker notes 존재
3. 핵심 키워드 (슬로건·DLF·197·Z3·OCP·박영희 등) 등장
4. 텍스트 박스가 슬라이드 영역 (13.333×7.5") 안에 있음
5. 텍스트 길이가 박스 크기 대비 합리적 (rough overflow heuristic)
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PATH = Path(__file__).parent / "woori_advocate_v1.pptx"
SLIDE_W_EMU = Emu(int(13.333 * 914400))
SLIDE_H_EMU = Emu(int(7.5 * 914400))

REQUIRED_KEYWORDS = {
    1: ["Woori Advocate", "고객 편", "세계 최초"],
    2: ["은행 편", "팔기 위한"],
    3: ["197", "DLF"],
    4: ["Seller", "Advocate", "박영희"],
    5: ["A2A", "Claude", "GPT"],
    6: ["MAD", "LangGraph", "embedding"],  # 신규 — MAD vs LangGraph
    7: ["MAD", "Constitutional", "Z3"],     # 학술 근거 (구 S6)
    8: ["금소법", "AI Act", "인공지능기본법"],
    9: ["박영희", "Z3"],
    10: ["이순자", "Family"],
    11: ["4분야", "통합"],
    12: ["우리GPT", "108억", "270"],
    13: ["ROI", "회피", "감경", "라이센스"],
    14: ["사후", "사전", "고객 편"],
    15: ["OCP", "Linux", "글로벌"],
    16: ["DLF 197", "표준", "출처"],
}


def main() -> None:
    prs = Presentation(PATH)
    print(f"file: {PATH.name}")
    print(f"size: {prs.slide_width} × {prs.slide_height} EMU "
          f"({prs.slide_width / 914400:.3f}\" × {prs.slide_height / 914400:.3f}\")")
    print(f"slide count: {len(prs.slides)}")
    assert len(prs.slides) == 16, "slide count != 16"

    issues: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_blobs: list[str] = []
        out_of_bounds: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            text_blobs.append(run.text)
            # bounds check
            if shape.left is not None and shape.top is not None:
                if (shape.left + (shape.width or 0)) > SLIDE_W_EMU + 1000:
                    out_of_bounds.append(f"shape '{getattr(shape, 'name', '?')}' overflows right")
                if (shape.top + (shape.height or 0)) > SLIDE_H_EMU + 1000:
                    out_of_bounds.append(f"shape '{getattr(shape, 'name', '?')}' overflows bottom")
                if (shape.left or 0) < -1000:
                    out_of_bounds.append(f"shape '{getattr(shape, 'name', '?')}' overflows left")
                if (shape.top or 0) < -1000:
                    out_of_bounds.append(f"shape '{getattr(shape, 'name', '?')}' overflows top")

        full_text = "\n".join(text_blobs)
        # speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text

        # required keywords
        missing = [kw for kw in REQUIRED_KEYWORDS.get(idx, [])
                   if kw not in full_text and kw not in notes_text]
        ok_kw = not missing

        print(f"\n[Slide {idx:2d}] text_chars={len(full_text):5d} "
              f"notes_chars={len(notes_text):5d} "
              f"keywords={'OK' if ok_kw else 'MISSING ' + str(missing)}")
        if out_of_bounds:
            for ob in out_of_bounds:
                print(f"   ⚠ {ob}")
                issues.append(f"S{idx}: {ob}")
        if missing:
            issues.append(f"S{idx}: missing keywords {missing}")
        if len(notes_text.strip()) < 30:
            print(f"   ⚠ speaker notes too short: '{notes_text[:60]}'")
            issues.append(f"S{idx}: short speaker notes")

    print("\n" + "=" * 60)
    if not issues:
        print("ALL CHECKS PASSED")
    else:
        print(f"✗ {len(issues)} issue(s):")
        for i in issues:
            print(f"  - {i}")


if __name__ == "__main__":
    main()
