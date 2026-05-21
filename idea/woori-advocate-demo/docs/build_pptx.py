"""
Woori Advocate 컨테스트 본선 PPT 생성기 (15장).

색상: Woori Navy (#003c69) 60% + Woori Blue (#0067ac) 20% + Advocate Red (#c8232b) 20%
폰트: 맑은 고딕 (한국어 최대 호환)
구조: Sandwich (S1·S15 dark, S2~S14 light)

실행:
    cd docs && python build_pptx.py
    → docs/woori_advocate_v1.pptx 생성
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


OUT = Path(__file__).parent / "woori_advocate_v1.pptx"

# ── Color palette ─────────────────────────────────────────────────────────
NAVY = RGBColor(0x00, 0x3C, 0x69)
BLUE = RGBColor(0x00, 0x67, 0xAC)
LIGHT_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
RED = RGBColor(0xC8, 0x23, 0x2B)
RED_DARK = RGBColor(0x8C, 0x10, 0x18)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF5, 0xF7, 0xFA)
INK = RGBColor(0x1B, 0x1F, 0x24)
MUTE = RGBColor(0x40, 0x49, 0x52)
BORDER = RGBColor(0xD6, 0xDD, 0xE4)

# 폰트
H_FONT = "맑은 고딕"
B_FONT = "맑은 고딕"


# ── 헬퍼 ──────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide, left, top, width, height, *,
    text: str = "",
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    font_name: str = B_FONT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(
    slide, left, top, width, height, items: list[tuple[str, int, bool, RGBColor]],
    align: int = PP_ALIGN.LEFT, anchor: int = MSO_ANCHOR.TOP,
    bg_color: RGBColor | None = None, border_color: RGBColor | None = None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.name = B_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    if bg_color is not None:
        # text box 자체는 fill 없으므로 사각형으로 대체할 일 있으면 add_card 사용
        pass
    return tb


def add_card(
    slide, left, top, width, height, *,
    fill: RGBColor = WHITE, border: RGBColor = BORDER,
    border_w: float = 1.0,
):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(border_w)
    shp.shadow.inherit = False
    return shp


def add_speaker_notes(slide, text: str) -> None:
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def add_footer(slide, idx: int, total: int = 16) -> None:
    """페이지 번호 + 브랜드 (모든 콘텐츠 슬라이드)"""
    add_text(
        slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.3),
        text=f"Woori Advocate · 우리은행 × SSAFY 2026",
        font_size=9, color=MUTE,
    )
    add_text(
        slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
        text=f"{idx} / {total}",
        font_size=9, color=MUTE, align=PP_ALIGN.RIGHT,
    )


# ── Presentation setup ───────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # blank


# ──────────────────────────────────────────────────────────────────────────
# S1 — Title + Slogan (dark navy)
# ──────────────────────────────────────────────────────────────────────────
def slide_01():
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)

    # 좌측 흰색 액센트 띠 대신 'W' 로고 마크 (좌상)
    mark = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.6), Inches(0.5), Inches(0.7), Inches(0.7))
    mark.adjustments[0] = 0.2
    mark.fill.solid()
    mark.fill.fore_color.rgb = WHITE
    mark.line.fill.background()
    tf = mark.text_frame
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "W"; r.font.name = H_FONT; r.font.bold = True
    r.font.size = Pt(28); r.font.color.rgb = NAVY

    add_text(s, Inches(1.4), Inches(0.55), Inches(6), Inches(0.6),
             text="우리은행", font_size=18, color=LIGHT_BLUE, font_name=H_FONT)

    # 메인 타이틀
    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
             text="Woori Advocate", font_size=66, bold=True, color=WHITE,
             font_name=H_FONT)
    add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(0.7),
             text="고객 편 AI · 사전 개입 · 이중 관점 동시 출력",
             font_size=22, color=LIGHT_BLUE, font_name=H_FONT)

    # 슬로건 카드
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.0))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x00, 0x2A, 0x4B)
    card.line.fill.background()
    add_text(s, Inches(0.9), Inches(4.45), Inches(11.5), Inches(0.8),
             text='"우리은행은 고객 편을 드는 AI를 둡니다"',
             font_size=30, bold=True, color=WHITE, font_name=H_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.6),
             text="세계 최초, 자기 행원에게 반대하는 AI",
             font_size=18, color=LIGHT_BLUE, font_name=H_FONT,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             text="우리은행 × SSAFY AI-금융소비자보호 아이디어 경진대회 · 2026.06.05 본선",
             font_size=12, color=LIGHT_BLUE, font_name=H_FONT,
             align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "안녕하십니까. 저희가 제안하는 시스템은 'Woori Advocate' 입니다. "
        "슬로건 그대로 — 우리은행은 고객 편을 드는 AI를 둡니다. "
        "세계 최초로, 자기 행원에게 반대하는 AI 입니다.\n\n"
        "[30초] 슬로건을 천천히, 또박또박. 청중이 슬로건을 머릿속에 새기게 함."
    )


# ──────────────────────────────────────────────────────────────────────────
# S2 — 문제 제기: 모든 은행 AI 는 '은행 편'
# ──────────────────────────────────────────────────────────────────────────
def slide_02():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="문제 제기", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="모든 은행 AI 는 \"은행 편\" 입니다",
             font_size=36, bold=True, color=NAVY, font_name=H_FONT)

    banks = [
        ("신한", "완전판매 AI\n스크립트·해피콜"),
        ("KB", "케이봇쌤\nLiiv AI 데스크"),
        ("하나", "하이(HAI)\n프로젝트 퍼스트"),
        ("우리", "AI-스미싱\n전기통신금융사기 모니터링"),
        ("JP Morgan", "IndexGPT\n셀링 보조"),
        ("Bank of\nAmerica", "Erica\nMerrill Insights"),
    ]
    # 2 rows × 3 cols on left half
    card_w = Inches(2.0); card_h = Inches(1.5); gap = Inches(0.15)
    for i, (name, desc) in enumerate(banks):
        col = i % 3; row = i // 3
        left = Inches(0.5) + (card_w + gap) * col
        top = Inches(2.2) + (card_h + gap) * row
        add_card(s, left, top, card_w, card_h, fill=SOFT, border=BORDER)
        add_text(s, left, top + Inches(0.15), card_w, Inches(0.5),
                 text=name, font_size=18, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, font_name=H_FONT)
        add_text(s, left, top + Inches(0.7), card_w, Inches(0.8),
                 text=desc, font_size=11, color=MUTE,
                 align=PP_ALIGN.CENTER, font_name=H_FONT)

    # 우측 질문 박스
    qbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(7.1), Inches(2.2), Inches(5.7), Inches(4.4))
    qbox.fill.solid(); qbox.fill.fore_color.rgb = NAVY
    qbox.line.fill.background()
    add_text(s, Inches(7.3), Inches(2.5), Inches(5.3), Inches(0.5),
             text="공통 라벨", font_size=12, color=LIGHT_BLUE,
             font_name=H_FONT)
    add_text(s, Inches(7.3), Inches(3.0), Inches(5.3), Inches(0.8),
             text="\"더 잘 팔기 위한 AI\"",
             font_size=22, bold=True, color=WHITE,
             font_name=H_FONT)
    add_text(s, Inches(7.3), Inches(4.0), Inches(5.3), Inches(0.4),
             text="——————", font_size=14, color=LIGHT_BLUE, font_name=H_FONT)
    add_text(s, Inches(7.3), Inches(4.4), Inches(5.3), Inches(1.8),
             text="그게 진정\n고객을 위한 것이었을까요?",
             font_size=26, bold=True, color=WHITE,
             font_name=H_FONT)

    add_footer(s, 2)
    add_speaker_notes(s,
        "지금 한국 4대 은행, 그리고 글로벌 JP Morgan, Bank of America 까지 — "
        "모든 은행 AI 는 공통 목표가 하나입니다. 상품을 더 잘 팔기 위해. "
        "신한의 완전판매 AI 스크립트, KB의 케이봇쌤, 하나의 AI 비서 — 전부 셀링 보조입니다. "
        "그것이 정말 고객을 위한 AI일까요?\n\n[45초] 'JP Morgan·BoA' 도 동일하다는 점에서 글로벌 차원의 빈틈."
    )


# ──────────────────────────────────────────────────────────────────────────
# S3 — DLF 사태: 197.1억
# ──────────────────────────────────────────────────────────────────────────
def slide_03():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="우리은행의 진짜 약점", font_size=14, color=RED, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="DLF 사태 — 은행권 역대 최대 과태료",
             font_size=32, bold=True, color=NAVY, font_name=H_FONT)

    # 큰 숫자 + 라벨 (좌)  — '억' '원' 분리로 오버플로 회피
    add_text(s, Inches(0.5), Inches(2.0), Inches(7), Inches(2.6),
             text="197.1억",
             font_size=110, bold=True, color=RED, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(4.5), Inches(7), Inches(0.5),
             text="원 · DLF 불완전판매 과태료 (2019)",
             font_size=16, bold=True, color=MUTE, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(5.0), Inches(7), Inches(0.4),
             text="은행권 역대 최대 규모",
             font_size=12, color=MUTE, font_name=H_FONT)

    # 보조 사실 (우)
    facts = [
        "사모펀드 신규판매 6개월 정지",
        "손태승 행장 문책경고 (중징계)",
        "분쟁조정 평균 배상비율 55%",
        "본질: 고객 프로필에 안 맞는 고위험 상품을 경고 없이 판매",
    ]
    add_card(s, Inches(7.6), Inches(2.2), Inches(5.2), Inches(2.7),
             fill=SOFT, border=BORDER)
    for i, line in enumerate(facts):
        add_text(s, Inches(7.85), Inches(2.4) + Inches(0.55) * i,
                 Inches(4.8), Inches(0.5),
                 text=f"▸ {line}", font_size=14, color=INK, font_name=H_FONT)

    # CTA 박스
    cta = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.3))
    cta.fill.solid(); cta.fill.fore_color.rgb = NAVY
    cta.line.fill.background()
    add_text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(0.5),
             text="만약 2019년에 Advocate 가 있었다면?",
             font_size=22, bold=True, color=WHITE, font_name=H_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.5),
             text="부적합 고객에게 가입 직전 5위험 표시 → 사태 자체가 없었음",
             font_size=14, color=LIGHT_BLUE, font_name=H_FONT,
             align=PP_ALIGN.CENTER)

    add_footer(s, 3)
    add_speaker_notes(s,
        "DLF 사태는 우리은행이 197억원, 은행권 역대 최대 과태료를 부담한 사건이었습니다. "
        "사모펀드 6개월 정지, 행장 문책경고. 본질은 단순합니다 — "
        "고객 프로필에 안 맞는 고위험 상품을 아무 경고 없이 판 것. "
        "만약 2019년에 Advocate 가 있었다면, 부적합 고객에게 가입 직전 5가지 위험을 알려줬을 것이고, "
        "이 사태 자체가 없었을 겁니다.\n\n[45초] 우리은행 본인의 약점을 정면으로 언급 — 신뢰성·진정성 확보."
    )


# ──────────────────────────────────────────────────────────────────────────
# S4 — Solution: 듀얼 모니터 (5초 침묵)
# ──────────────────────────────────────────────────────────────────────────
def slide_04():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="해결책", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.0),
             text="가입 직전, 같은 화면을 둘로 가른다",
             font_size=30, bold=True, color=NAVY, font_name=H_FONT)

    # 좌: Seller pane (navy)
    seller = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.5), Inches(2.0), Inches(6.1), Inches(4.5))
    seller.fill.solid(); seller.fill.fore_color.rgb = NAVY
    seller.line.fill.background()
    add_text(s, Inches(0.7), Inches(2.1), Inches(5.7), Inches(0.5),
             text="행원 PC — Seller 관점", font_size=14, bold=True, color=WHITE,
             font_name=H_FONT)
    add_text(s, Inches(0.7), Inches(2.55), Inches(5.7), Inches(0.4),
             text="Claude · 셀링 포인트", font_size=10,
             color=LIGHT_BLUE, font_name=H_FONT)

    seller_lines = [
        "▸ 시중 예적금 대비 약 3배,",
        "   연 최대 6.0% 세전 쿠폰 예상",
        "▸ 적극투자형 이상 적합 상품",
        "▸ 6개월마다 조기상환 평가",
        "",
        "[Z3 SMT ✓] 부당권유 패턴 0개",
    ]
    for i, t in enumerate(seller_lines):
        add_text(s, Inches(0.85), Inches(3.1) + Inches(0.42) * i,
                 Inches(5.5), Inches(0.4),
                 text=t, font_size=13, color=WHITE, font_name=H_FONT)

    # 우: Advocate pane (red)
    adv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.5))
    adv.fill.solid(); adv.fill.fore_color.rgb = RED
    adv.line.fill.background()
    add_text(s, Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.5),
             text="고객 태블릿 — Advocate 관점", font_size=14, bold=True,
             color=WHITE, font_name=H_FONT)
    add_text(s, Inches(7.0), Inches(2.55), Inches(5.6), Inches(0.4),
             text="OpenAI GPT · 5위험 + 출처 표기", font_size=10,
             color=LIGHT_BLUE, font_name=H_FONT)

    adv_lines = [
        "1. 3년간 자금 묶임  [출처: T-002]",
        "   의료비 1,200만 지출 충돌",
        "2. 5년간 고위험 거래 0건 [T-004]",
        "3. 원금 최대 100% 손실 [T-001]",
        "4. 유사 상품 손실 7.2% [T-006]",
        "5. 75세 + 변동성 부담 [T-005]",
    ]
    for i, t in enumerate(adv_lines):
        add_text(s, Inches(7.15), Inches(3.1) + Inches(0.42) * i,
                 Inches(5.6), Inches(0.4),
                 text=t, font_size=13, color=WHITE, font_name=H_FONT)

    # 하단 컨텍스트
    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
             text="박영희 75세 · 의료비 1,200만 지출 · 고위험 상품 거래 5년 0건 · 자칭 손주 압박 전화 24시간 내",
             font_size=11, color=MUTE, font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 4)
    add_speaker_notes(s,
        "저희의 해결책은 이렇습니다.\n\n"
        "[★★ 5초 침묵 — 정확히 5초. 3초 짧고 7초 어색. 화면을 그대로 보여줌 ★★]\n\n"
        "보시는 화면이 가입 직전 영업점입니다. "
        "행원의 모니터는 상품 강점을 말하고, 고객 앞 태블릿은 같은 순간 "
        "이 상품이 맞지 않는 이유를 말합니다. "
        "좌측은 Seller, 우측은 Advocate — "
        "같은 LLM을 두 페르소나로 격리시켜 이중 관점을 동시에 출력합니다."
    )


# ──────────────────────────────────────────────────────────────────────────
# S5 — MAD Architecture
# ──────────────────────────────────────────────────────────────────────────
def slide_05():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="아키텍처", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="MAD + 이질 LLM ensemble + A2A 표준",
             font_size=30, bold=True, color=NAVY, font_name=H_FONT)

    # Orchestrator (center top)
    orch = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(4.7), Inches(2.0), Inches(4.0), Inches(0.9))
    orch.adjustments[0] = 0.2
    orch.fill.solid(); orch.fill.fore_color.rgb = NAVY
    orch.line.fill.background()
    add_text(s, Inches(4.7), Inches(2.05), Inches(4.0), Inches(0.4),
             text="Advocate Orchestrator", font_size=15, bold=True,
             color=WHITE, font_name=H_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.7), Inches(2.45), Inches(4.0), Inches(0.4),
             text=".well-known/agent.json · A2A 표준",
             font_size=10, color=LIGHT_BLUE, font_name=H_FONT,
             align=PP_ALIGN.CENTER)

    # 8 agent cards (2 rows × 4 cols) below
    agents = [
        ("Seller", "Claude", BLUE),
        ("Advocate", "GPT", RED),
        ("Regulator", "Llama + Z3", NAVY),
        ("History", "거래·민원", MUTE),
        ("Peer", "통계 mock", MUTE),
        ("Senior", "65+ 자동", BLUE),
        ("FDS", "사기 패턴", RED),
        ("Family", "동의 절차", BLUE),
    ]
    card_w = Inches(2.8); card_h = Inches(1.3); gap_x = Inches(0.2); gap_y = Inches(0.2)
    grid_left = Inches(0.5)
    for i, (name, sub, accent) in enumerate(agents):
        col = i % 4; row = i // 4
        left = grid_left + (card_w + gap_x) * col
        top = Inches(3.4) + (card_h + gap_y) * row
        add_card(s, left, top, card_w, card_h, fill=WHITE, border=accent, border_w=2.0)
        add_text(s, left, top + Inches(0.18), card_w, Inches(0.5),
                 text=name, font_size=16, bold=True, color=accent,
                 align=PP_ALIGN.CENTER, font_name=H_FONT)
        add_text(s, left, top + Inches(0.68), card_w, Inches(0.5),
                 text=sub, font_size=11, color=MUTE,
                 align=PP_ALIGN.CENTER, font_name=H_FONT)

    # 하단 라벨
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
             text="이질 LLM 5종 ensemble (Claude · GPT · Llama · Qwen · Solar) · A2A 메시지 → OCP audit 5년 보관",
             font_size=12, color=MUTE, font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 5)
    add_speaker_notes(s,
        "내부 구조입니다. Multi-Agent Debate 아키텍처 위에 이질 LLM ensemble — "
        "Seller 는 Anthropic Claude, Advocate 는 OpenAI GPT, Regulator 는 Llama + Z3 형식 검증을 결합합니다. "
        "서로 다른 학습 데이터, 다른 정렬 철학을 가진 모델들이 진짜 대립합니다. "
        "8개 에이전트는 Google A2A 표준으로 통신해서 우리은행의 기존 175개 AI 에이전트와 "
        "즉시 상호운용 가능합니다.\n\n[60초] '이질 LLM' 과 'A2A' — 두 키워드 강조."
    )


# ──────────────────────────────────────────────────────────────────────────
# S6 — MAD vs LangGraph 비교 (신규 아키텍처 슬라이드)
# ──────────────────────────────────────────────────────────────────────────
def slide_06_mad_vs_langgraph():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="왜 단순 LangGraph 가 아닌가",
             font_size=14, color=RED, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.0),
             text="MAD (Multi-Agent Debate) — cross-context embedding",
             font_size=28, bold=True, color=NAVY, font_name=H_FONT)

    # 좌: LangGraph 다이어그램 + 한계
    add_card(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.5),
             fill=SOFT, border=BORDER, border_w=1.5)
    add_text(s, Inches(0.7), Inches(2.15), Inches(5.5), Inches(0.5),
             text="LangGraph (상태 그래프)", font_size=16, bold=True,
             color=MUTE, font_name=H_FONT)
    add_text(s, Inches(0.7), Inches(2.55), Inches(5.5), Inches(0.4),
             text="단일 LLM 이 노드 별로 도구·라우팅 수행",
             font_size=11, color=MUTE, font_name=H_FONT)

    # LangGraph 미니 다이어그램 (텍스트 박스 연결)
    ln_nodes = [
        ("Node A", Inches(0.85), Inches(3.1)),
        ("Node B", Inches(2.6), Inches(3.1)),
        ("Node C", Inches(4.35), Inches(3.1)),
    ]
    for name, left, top in ln_nodes:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, Inches(1.5), Inches(0.6))
        box.adjustments[0] = 0.2
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = MUTE; box.line.width = Pt(1.5)
        add_text(s, left, top + Inches(0.13), Inches(1.5), Inches(0.4),
                 text=name, font_size=11, bold=True, color=MUTE,
                 align=PP_ALIGN.CENTER, font_name=H_FONT)

    # 한계 리스트
    limits = [
        "× 단일 LLM — 다양성 0",
        "× 노드 간 메시지는 state 전달만",
        "× '진짜 대립' 메커니즘 부재",
        "× 환각이 그래프 전체에 전파",
        "× 학술 근거 없음 (단순 프레임워크)",
    ]
    for i, line in enumerate(limits):
        add_text(s, Inches(0.85), Inches(4.0) + Inches(0.4) * i,
                 Inches(5.2), Inches(0.4),
                 text=line, font_size=12, color=INK, font_name=H_FONT)

    # 우: MAD 다이어그램 + 이점
    add_card(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.5),
             fill=WHITE, border=RED, border_w=2.5)
    add_text(s, Inches(7.0), Inches(2.15), Inches(5.5), Inches(0.5),
             text="우리 MAD (이질 LLM 토론)", font_size=16, bold=True,
             color=RED, font_name=H_FONT)
    add_text(s, Inches(7.0), Inches(2.55), Inches(5.5), Inches(0.4),
             text="라운드 N → 상대 출력이 내 system prompt 에 embedding",
             font_size=11, color=MUTE, font_name=H_FONT)

    # MAD 미니 다이어그램: 2 LLM 박스 + 양방향 화살표
    seller_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(7.0), Inches(3.1),
                                    Inches(2.3), Inches(0.7))
    seller_box.adjustments[0] = 0.15
    seller_box.fill.solid(); seller_box.fill.fore_color.rgb = BLUE
    seller_box.line.fill.background()
    add_text(s, Inches(7.0), Inches(3.18), Inches(2.3), Inches(0.4),
             text="Seller", font_size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name=H_FONT)
    add_text(s, Inches(7.0), Inches(3.55), Inches(2.3), Inches(0.3),
             text="Claude (Anthropic)", font_size=10, color=LIGHT_BLUE,
             align=PP_ALIGN.CENTER, font_name=H_FONT)

    adv_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(10.3), Inches(3.1),
                                 Inches(2.3), Inches(0.7))
    adv_box.adjustments[0] = 0.15
    adv_box.fill.solid(); adv_box.fill.fore_color.rgb = RED
    adv_box.line.fill.background()
    add_text(s, Inches(10.3), Inches(3.18), Inches(2.3), Inches(0.4),
             text="Advocate", font_size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name=H_FONT)
    add_text(s, Inches(10.3), Inches(3.55), Inches(2.3), Inches(0.3),
             text="GPT (OpenAI)", font_size=10, color=LIGHT_BLUE,
             align=PP_ALIGN.CENTER, font_name=H_FONT)

    # 양방향 화살표 (라운드 2 embedding)
    ar_lr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(9.35), Inches(3.3),
                               Inches(0.85), Inches(0.25))
    ar_lr.fill.solid(); ar_lr.fill.fore_color.rgb = RED_DARK
    ar_lr.line.fill.background()

    ar_rl = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW,
                               Inches(9.35), Inches(3.65),
                               Inches(0.85), Inches(0.25))
    ar_rl.fill.solid(); ar_rl.fill.fore_color.rgb = RED_DARK
    ar_rl.line.fill.background()

    # 이점 리스트
    benefits = [
        "✓ 이질 LLM (Anthropic + OpenAI) — 학습 데이터·정렬 다름",
        "✓ 상대 출력을 context 에 embedding → 진짜 대립",
        "✓ 라운드 반복 → Moderator 합의/불일치 추출",
        "✓ MIT MAD ICML 2024 +80.9% 정확도 literal 적용",
        "✓ 환각: cross-LLM 검증으로 < 0.5% 목표",
    ]
    for i, line in enumerate(benefits):
        add_text(s, Inches(7.0), Inches(4.4) + Inches(0.4) * i,
                 Inches(5.8), Inches(0.4),
                 text=line, font_size=12, color=INK, font_name=H_FONT, bold=True)

    # 하단 결론
    cta = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5))
    cta.fill.solid(); cta.fill.fore_color.rgb = NAVY
    cta.line.fill.background()
    add_text(s, Inches(0.7), Inches(6.78), Inches(11.9), Inches(0.4),
             text='Q5 ("같은 LLM 진짜 대립?") → 답: 다른 LLM, 다른 출처 → 단순 LangGraph 가 줄 수 없는 학술적 근거',
             font_size=12, color=WHITE, bold=True, font_name=H_FONT,
             align=PP_ALIGN.CENTER)

    add_footer(s, 6, total=16)
    add_speaker_notes(s,
        "왜 단순 LangGraph 가 아닌 MAD 인가. "
        "LangGraph 는 좋은 상태 그래프 프레임워크지만 — 단일 LLM 이 노드를 옮겨 다닐 뿐, "
        "에이전트 간 '진짜 대립' 이 없습니다. "
        "우리 MAD 는 다릅니다. 라운드 1 은 독립 호출, 라운드 2 부터 "
        "상대의 직전 출력이 내 system prompt context 에 embedding 됩니다. "
        "그래서 Seller(Claude) 와 Advocate(GPT) 가 서로의 주장을 보고 반응합니다. "
        "MIT MAD 논문의 +80.9% 정확도 향상이 이 cross-context embedding 에서 나옵니다.\n\n"
        "[45초] 'LangGraph 도 좋지만 우리는 한 단계 더 갔다' 톤으로."
    )


# ──────────────────────────────────────────────────────────────────────────
# S7 — 학술 근거 (기존 S6 자리)
# ──────────────────────────────────────────────────────────────────────────
def slide_06():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="학술 근거", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="추측이 아닌 논문 기반",
             font_size=30, bold=True, color=NAVY, font_name=H_FONT)

    papers = [
        ("MIT MAD", "ICML 2024", "여러 LLM 토론 시\n사실 정확도 +80.9%\n환각 ↓", BLUE),
        ("Constitutional AI", "Anthropic 2022", "헌법 자가 검증 →\n금소법 6대 의무를\nAdvocate 헌법에 삽입", NAVY),
        ("System 1 / 2", "Kahneman 2011", "Seller = 직관(S1)\nAdvocate = 분석(S2)\n오류율 ↓", BLUE),
        ("Z3 한국어 금융 SMT", "AAAI/KDD 2026 publishable", "한국어 + 금융 + LLM\n출력 형식 검증\n학술 첫 적용", RED),
    ]
    card_w = Inches(2.95); card_h = Inches(3.5); gap = Inches(0.15)
    for i, (t, sub, body, accent) in enumerate(papers):
        left = Inches(0.5) + (card_w + gap) * i
        top = Inches(2.3)
        add_card(s, left, top, card_w, card_h, fill=WHITE, border=accent, border_w=2.5)
        # 상단 액센트 박스
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, card_w, Inches(0.7))
        head.fill.solid(); head.fill.fore_color.rgb = accent
        head.line.fill.background()
        add_text(s, left, top + Inches(0.05), card_w, Inches(0.6),
                 text=t, font_size=16, bold=True, color=WHITE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, left, top + Inches(0.85), card_w, Inches(0.5),
                 text=sub, font_size=11, color=MUTE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.15), top + Inches(1.5),
                 card_w - Inches(0.3), Inches(1.8),
                 text=body, font_size=13, color=INK,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             text="\"+80.9%\" 는 비유가 아니라 우리 ensemble 구성에 그대로 literal 적용",
             font_size=12, color=NAVY, bold=True,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 6)
    add_speaker_notes(s,
        "이 아키텍처는 추측이 아니라 학술 근거를 갖습니다. "
        "MIT MAD 논문은 여러 LLM 이 토론하면 단일 LLM 대비 사실 정확도가 +80.9% 향상되고 "
        "환각이 감소함을 증명했습니다. "
        "Anthropic Constitutional AI 는 헌법 원칙 자가 검증의 표준을 만들었고, "
        "저희는 그 헌법에 금소법 6대 의무를 직접 삽입했습니다. "
        "그리고 한국어 금융 도메인에 Z3 SMT 형식 검증을 적용한 것은 학술적으로 첫 사례입니다.\n\n"
        "[45초] '+80.9%' 와 '학술 첫 적용' — 두 숫자."
    )


# ──────────────────────────────────────────────────────────────────────────
# S7 — 규제 정합성
# ──────────────────────────────────────────────────────────────────────────
def slide_07():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="규제 정합성", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="2026년 — 한·EU 동시 시행기",
             font_size=30, bold=True, color=NAVY, font_name=H_FONT)

    regs = [
        ("한국 금소법", "2021.03 시행\n2026.01 감독규정 개정",
         "적합성·적정성\n설명·불공정\n부당권유·광고",
         "Z3 §21·§22\nMod 적합성 검증"),
        ("인공지능기본법", "2026.01 시행",
         "고위험 AI 윤리위\n사후 검증 의무",
         "OCP audit 5년\nAgent Card 공개"),
        ("EU AI Act", "2026.08 본격 집행",
         "투명성 / 인간감독\n5년 기록 / 매출 7% 과징금",
         "OCP envelope_digest\n자동 충족"),
    ]
    card_w = Inches(4.0); card_h = Inches(4.4); gap = Inches(0.2)
    for i, (name, when, what, ours) in enumerate(regs):
        left = Inches(0.5) + (card_w + gap) * i
        top = Inches(2.2)
        add_card(s, left, top, card_w, card_h, fill=WHITE, border=NAVY, border_w=1.5)
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, card_w, Inches(0.8))
        head.fill.solid(); head.fill.fore_color.rgb = NAVY
        head.line.fill.background()
        add_text(s, left, top + Inches(0.1), card_w, Inches(0.6),
                 text=name, font_size=18, bold=True, color=WHITE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, left, top + Inches(0.95), card_w, Inches(0.5),
                 text=when, font_size=11, color=MUTE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, left + Inches(0.2), top + Inches(1.6),
                 card_w - Inches(0.4), Inches(1.3),
                 text="법 요구:", font_size=11, color=MUTE, font_name=H_FONT, bold=True)
        add_text(s, left + Inches(0.2), top + Inches(1.9),
                 card_w - Inches(0.4), Inches(1.3),
                 text=what, font_size=13, color=INK, font_name=H_FONT)
        add_text(s, left + Inches(0.2), top + Inches(3.05),
                 card_w - Inches(0.4), Inches(0.4),
                 text="Advocate 충족:", font_size=11, color=RED,
                 font_name=H_FONT, bold=True)
        add_text(s, left + Inches(0.2), top + Inches(3.4),
                 card_w - Inches(0.4), Inches(1.0),
                 text=ours, font_size=13, color=INK, font_name=H_FONT, bold=True)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             text="Advocate 도입 = 3법 자동 충족 + EU 매출 7% 과징금 사전 차단",
             font_size=12, color=NAVY, bold=True,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 7)
    add_speaker_notes(s,
        "규제 측면입니다. 2026년은 한국 인공지능기본법, EU AI Act 가 본격 작동하는 해입니다. "
        "Woori Advocate 는 OCP audit log 와 Z3 형식 증명으로 3개 법규를 한 번에 자동 충족 시킵니다. "
        "EU AI Act 위반 시 매출의 7% 과징금을 사전 차단합니다. "
        "우리은행이 컴플라이언스 표준 reference 은행이 되는 길입니다.\n\n[45초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S8 — 시나리오 1 ELS 시연 (메인)
# ──────────────────────────────────────────────────────────────────────────
def slide_08():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="LIVE DEMO", font_size=14, color=RED, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="시나리오 1 — 박영희 75세 ELS 가입",
             font_size=30, bold=True, color=NAVY, font_name=H_FONT)

    # 좌: 5 단계 flow
    steps = [
        ("01", "행원 PC 에 Seller 자동 표시"),
        ("02", "고객 태블릿 시니어 모드 활성"),
        ("03", "Advocate 5위험 + 출처 표기"),
        ("04", "Z3 SMT 검증 배지"),
        ("05", "5개 체크박스 + 자필 서명"),
    ]
    add_text(s, Inches(0.5), Inches(2.0), Inches(6), Inches(0.5),
             text="시연 흐름", font_size=14, bold=True, color=BLUE, font_name=H_FONT)
    for i, (n, t) in enumerate(steps):
        top = Inches(2.55) + Inches(0.75) * i
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.5), top, Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()
        add_text(s, Inches(0.5), top + Inches(0.05), Inches(0.55), Inches(0.45),
                 text=n, font_size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name=H_FONT)
        add_text(s, Inches(1.2), top + Inches(0.1), Inches(5.0), Inches(0.5),
                 text=t, font_size=14, color=INK, font_name=H_FONT)

    # 우: Demo preview card
    pv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(6.7), Inches(2.0), Inches(6.1), Inches(4.5))
    pv.fill.solid(); pv.fill.fore_color.rgb = SOFT
    pv.line.color.rgb = BORDER
    add_text(s, Inches(6.9), Inches(2.15), Inches(5.7), Inches(0.4),
             text="브라우저: http://localhost:8000", font_size=11,
             color=MUTE, font_name=H_FONT)
    # 미니 좌우 분할
    mini_l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(6.95), Inches(2.7), Inches(2.85), Inches(3.5))
    mini_l.fill.solid(); mini_l.fill.fore_color.rgb = NAVY
    mini_l.line.fill.background()
    add_text(s, Inches(7.05), Inches(2.8), Inches(2.6), Inches(0.4),
             text="Seller", font_size=12, bold=True, color=WHITE, font_name=H_FONT)
    for i, t in enumerate(["▸ 연 최대 6% 쿠폰",
                           "▸ 적극투자형 적합",
                           "▸ 6개월 평가",
                           "✓ Z3 통과"]):
        add_text(s, Inches(7.05), Inches(3.25) + Inches(0.5) * i,
                 Inches(2.6), Inches(0.4),
                 text=t, font_size=11, color=WHITE, font_name=H_FONT)

    mini_r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(9.85), Inches(2.7), Inches(2.85), Inches(3.5))
    mini_r.fill.solid(); mini_r.fill.fore_color.rgb = RED
    mini_r.line.fill.background()
    add_text(s, Inches(9.95), Inches(2.8), Inches(2.6), Inches(0.4),
             text="Advocate", font_size=12, bold=True, color=WHITE, font_name=H_FONT)
    for i, t in enumerate(["1. 자금 묶임 [T-002]",
                           "2. 거래 0건 [T-004]",
                           "3. 원금 100% [T-001]",
                           "4. 손실 7.2% [T-006]",
                           "5. 변동성 [T-005]"]):
        add_text(s, Inches(9.95), Inches(3.25) + Inches(0.45) * i,
                 Inches(2.6), Inches(0.4),
                 text=t, font_size=10, color=WHITE, font_name=H_FONT)

    add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             text="목표 응답 시간 < 2.5초 p95 · 캐시 hit 시 < 500ms",
             font_size=11, color=MUTE, font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 8)
    add_speaker_notes(s,
        "실시간 시연입니다.\n"
        "[클릭] '박영희 ELS 가입 시도 → Advocate 발동' 버튼 클릭.\n"
        "박영희 75세 고객이 홍콩H지수 연계 ELS 에 가입을 시도합니다. 화면이 분할됩니다.\n\n"
        "좌측 행원 PC — Seller 가 셀링 포인트 4개를 보여주고, "
        "Z3 SMT 검증이 통과됩니다 — 부당권유 패턴 0개, 수학적 증명.\n\n"
        "우측 고객 태블릿 — Advocate 가 5가지 위험을 박영희 본인 데이터와 매핑해서 보여줍니다. "
        "각 항목 끝에 약관 청크 ID — 환각 차단. 박영희님의 의료비 1,200만 지출, "
        "5년간 고위험 상품 거래 0건이 위험 1·2번에 반영됩니다.\n\n"
        "고객은 5개 체크박스를 모두 누른 뒤에만 가입 진행이 가능합니다. "
        "체크 안 누르면 행원 호출.\n\n[90초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S9 — 시나리오 2: 시니어 + 포용금융
# ──────────────────────────────────────────────────────────────────────────
def slide_09():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="시나리오 2 — 시니어 + 포용금융",
             font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="이순자 75세, 손주 명의 1,500만 송금",
             font_size=28, bold=True, color=NAVY, font_name=H_FONT)

    # 좌: 시니어 UI mockup
    phone = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.7), Inches(2.0), Inches(4.0), Inches(5.0))
    phone.adjustments[0] = 0.08
    phone.fill.solid(); phone.fill.fore_color.rgb = RED
    phone.line.fill.background()
    add_text(s, Inches(0.85), Inches(2.15), Inches(3.7), Inches(0.5),
             text="시니어 모드 (자동)", font_size=12, bold=True, color=WHITE,
             font_name=H_FONT)
    add_text(s, Inches(0.85), Inches(2.6), Inches(3.7), Inches(0.6),
             text="이순자 님",
             font_size=22, bold=True, color=WHITE, font_name=H_FONT)
    add_text(s, Inches(0.85), Inches(3.2), Inches(3.7), Inches(0.6),
             text="송금 전 한 번만\n확인해 주세요",
             font_size=18, color=WHITE, font_name=H_FONT)
    add_text(s, Inches(0.85), Inches(4.4), Inches(3.7), Inches(2.0),
             text="✓ 큰 글씨\n✓ TTS 음성 안내\n✓ 한 번에 한 카드씩\n✓ 가족 동의 푸시",
             font_size=15, color=WHITE, font_name=H_FONT)

    # 우: FDS + Family 흐름
    add_text(s, Inches(5.2), Inches(2.0), Inches(7.5), Inches(0.5),
             text="FDS 4개 시그널 검출", font_size=16, bold=True, color=RED,
             font_name=H_FONT)
    signals = [
        ("①", "비정상 거액: 평소 평균 이체액의 50배"),
        ("②", "수취 계좌 최근 개설 (3일 전)"),
        ("③", "자칭 '손주' 와 최근 6개월 통화 0회"),
        ("④", "24시간 내 의심 발신 통화 24분"),
    ]
    for i, (n, t) in enumerate(signals):
        top = Inches(2.55) + Inches(0.6) * i
        add_text(s, Inches(5.2), top, Inches(0.5), Inches(0.5),
                 text=n, font_size=18, bold=True, color=RED, font_name=H_FONT)
        add_text(s, Inches(5.7), top + Inches(0.05), Inches(7), Inches(0.5),
                 text=t, font_size=13, color=INK, font_name=H_FONT)

    add_text(s, Inches(5.2), Inches(5.1), Inches(7.5), Inches(0.5),
             text="Family Consent 절차 발동", font_size=16, bold=True, color=NAVY,
             font_name=H_FONT)
    flow = [
        "1. 등록 손주 직접 통화 권유",
        "2. 자녀 안심앱 푸시 알림 동시 발송",
        "3. 30분 미응답 시 안심센터 자동 연결",
    ]
    for i, t in enumerate(flow):
        add_text(s, Inches(5.2), Inches(5.55) + Inches(0.4) * i,
                 Inches(7.5), Inches(0.4),
                 text=t, font_size=12, color=INK, font_name=H_FONT)

    # 응모 분야 매핑
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             text="응모분야 1 (취약계층 보호) + 4 (보이스피싱 예방) 동시 처리",
             font_size=12, color=NAVY, bold=True,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 9)
    add_speaker_notes(s,
        "시니어 시나리오입니다. 65세 이상 고객은 자동으로 시니어 모드 — "
        "큰 글씨, TTS 음성 안내, 한 번에 한 위험 항목씩. "
        "손주 명의 계좌로 1,500만원 송금을 시도하면 FDS 가 4개 시그널을 검출합니다 — "
        "평소 50배 거액, 신규 계좌 3일, 통화 이력 0회, 의심 발신 통화. "
        "Family Consent 절차가 발동되어 진짜 손주에게 직접 전화 확인을 권유합니다. "
        "응모분야 1번 취약계층 보호와 4번 보이스피싱을 동시에 처리합니다.\n\n[45초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S10 — 통합 시나리오 0
# ──────────────────────────────────────────────────────────────────────────
def slide_10():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="통합 시나리오 0", font_size=14, color=RED, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="박영희 75세 + 손주 사칭 ELS — 4분야 동시 발동",
             font_size=26, bold=True, color=NAVY, font_name=H_FONT)

    # 2x2 grid — 4 응모분야
    quads = [
        ("응모분야 1", "취약계층 보호", "Senior + Family Consent",
         "75세 + senior_mode 플래그\n자동 시니어 UI · TTS", BLUE),
        ("응모분야 2", "불완전판매 예방", "Advocate + Z3 + Moderator",
         "적합성 위반 5위험\nZ3 형식 증명 · 청크 출처", NAVY),
        ("응모분야 3", "민원 사전 예방", "감정 분석 + 행원 알람",
         "'압박감' 패턴 검출\n행원 PC 우상단 알람", RED),
        ("응모분야 4", "보이스피싱 예방", "FDS + Family Consent",
         "신규 계좌 + 통화 0회\n+ 평소 50배 거액", RED_DARK),
    ]
    card_w = Inches(6.0); card_h = Inches(2.4); gap = Inches(0.2)
    for i, (tag, name, agents, body, accent) in enumerate(quads):
        col = i % 2; row = i // 2
        left = Inches(0.5) + (card_w + gap) * col
        top = Inches(2.1) + (card_h + gap) * row
        add_card(s, left, top, card_w, card_h, fill=WHITE, border=accent, border_w=2.0)
        # 좌측 색 띠
        side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, Inches(0.25), card_h)
        side.fill.solid(); side.fill.fore_color.rgb = accent
        side.line.fill.background()

        add_text(s, left + Inches(0.4), top + Inches(0.15),
                 Inches(2.5), Inches(0.4),
                 text=tag, font_size=11, color=accent, bold=True, font_name=H_FONT)
        add_text(s, left + Inches(0.4), top + Inches(0.5),
                 Inches(5.5), Inches(0.5),
                 text=name, font_size=18, bold=True, color=NAVY, font_name=H_FONT)
        add_text(s, left + Inches(0.4), top + Inches(1.0),
                 Inches(5.5), Inches(0.4),
                 text=agents, font_size=11, color=MUTE, font_name=H_FONT)
        add_text(s, left + Inches(0.4), top + Inches(1.4),
                 Inches(5.5), Inches(0.9),
                 text=body, font_size=12, color=INK, font_name=H_FONT)

    add_text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
             text="1개 시나리오 = 4개 응모분야 모두 발동 → 효과성 만점 논리",
             font_size=12, color=NAVY, bold=True,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 10)
    add_speaker_notes(s,
        "더 강력한 시나리오를 보여드립니다. 박영희 75세, 자칭 손주에게 사기 전화를 받은 뒤 "
        "영업점에서 ELS 가입을 시도합니다.\n\n"
        "동시에 4개 응모분야가 모두 발동합니다. "
        "(1) 시니어 모드 자동 활성. "
        "(2) Advocate 가 적합성 위반 5가지 위험 표시 + Z3 검증. "
        "(3) 가입 후 손주 계좌 송금 단계에서 FDS 가 4개 시그널 검출. "
        "(4) 박영희 음성에서 '압박감' 감정이 감지되어 행원 PC 에 사전 알람. "
        "단 하나의 시나리오로, 응모 4분야 전체를 커버합니다.\n\n[90초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S11 — 실현가능성: 자산 + 비용 분담
# ──────────────────────────────────────────────────────────────────────────
def slide_11():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="실현가능성", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="우리은행 보유 자산 + 4-way 비용 분담",
             font_size=28, bold=True, color=NAVY, font_name=H_FONT)

    # 좌: 자산 5개 (icon + label)
    add_text(s, Inches(0.5), Inches(2.0), Inches(6), Inches(0.5),
             text="우리은행 보유 자산 (개발 비용 거의 0)",
             font_size=14, bold=True, color=NAVY, font_name=H_FONT)
    assets = [
        ("우리GPT", "2024.12 · 1,000만건 내부 데이터", "Seller·Advocate 베이스"),
        ("운영GPT", "2025.05 · 리스크 특화", "Regulator·Moderator 재활용"),
        ("175 AI 에이전트", "2026.03 배치", "A2A 표준으로 즉시 통합"),
        ("GPU 클러스터", "2026.04 자체 구축", "온프레미스 추론"),
        ("STT/AICC", "영업점 태블릿 보유", "시나리오 즉시 가동"),
    ]
    for i, (name, when, use) in enumerate(assets):
        top = Inches(2.55) + Inches(0.78) * i
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(0.5), top + Inches(0.1),
                                 Inches(0.35), Inches(0.35))
        dot.fill.solid(); dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_text(s, Inches(0.5), top + Inches(0.13), Inches(0.35), Inches(0.3),
                 text=str(i+1), font_size=12, bold=True, color=WHITE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.0), top, Inches(5.5), Inches(0.4),
                 text=name, font_size=14, bold=True, color=NAVY, font_name=H_FONT)
        add_text(s, Inches(1.0), top + Inches(0.4), Inches(5.5), Inches(0.4),
                 text=f"{when}  →  {use}",
                 font_size=10, color=MUTE, font_name=H_FONT)

    # 우: 비용 분담 (도넛 대신 4-bar)
    add_text(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(0.5),
             text="P0~P3 누적 270억 비용 분담",
             font_size=14, bold=True, color=NAVY, font_name=H_FONT)
    splits = [
        ("우리은행 자체", 40, "108억", BLUE),
        ("금융위 디지털금융 R&D", 30, "81억", NAVY),
        ("과기정통부 표준화", 15, "40.5억", RED),
        ("KAIST·LF·SSAFY", 15, "40.5억", MUTE),
    ]
    bar_full = 5.0
    for i, (name, pct, amt, color) in enumerate(splits):
        top = Inches(2.7) + Inches(0.95) * i
        add_text(s, Inches(7.0), top, Inches(4), Inches(0.4),
                 text=name, font_size=12, bold=True, color=INK, font_name=H_FONT)
        add_text(s, Inches(11.0), top, Inches(1.8), Inches(0.4),
                 text=f"{pct}% · {amt}", font_size=12, bold=True,
                 color=color, font_name=H_FONT, align=PP_ALIGN.RIGHT)
        # bar bg
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(7.0), top + Inches(0.45),
                                Inches(bar_full), Inches(0.25))
        bg.fill.solid(); bg.fill.fore_color.rgb = SOFT
        bg.line.fill.background()
        # bar fg
        fg_w = bar_full * pct / 50  # 50% 가 풀바 절반보다 약간 크게
        fg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(7.0), top + Inches(0.45),
                                Inches(min(fg_w, bar_full)), Inches(0.25))
        fg.fill.solid(); fg.fill.fore_color.rgb = color
        fg.line.fill.background()

    add_text(s, Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.4),
             text="P4·P5 (200억) 은 라이센스 수익 선납으로 자력 가능",
             font_size=11, color=MUTE, font_name=H_FONT)

    add_footer(s, 11)
    add_speaker_notes(s,
        "실현가능성. 우리은행이 이미 갖고 있는 5개 자산 — "
        "우리GPT, 운영GPT, 175개 AI 에이전트, GPU 클러스터, 영업점 STT — 위에 얹습니다. "
        "외부 인프라 신규 투자 거의 0. "
        "그리고 비용은 우리은행 단독이 아니라 4-way 분담 — "
        "인공지능기본법 시범사업으로 금융위 R&D 30%, 표준화 과제로 과기정통부 15%, "
        "학술·표준 컨소시엄 15%. P0~P3 누적 270억 중 우리은행 자체는 108억.\n\n[45초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S12 — ROI 3축
# ──────────────────────────────────────────────────────────────────────────
def slide_12():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="ROI", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="3축 수익 모델 · P1 6개월 종료 시 Break-even",
             font_size=26, bold=True, color=NAVY, font_name=H_FONT)

    axes = [
        ("A", "회피", "DLF·홍콩 ELS 급\n사태 1건 회피",
         "197억", "1회만으로 P1 비용 회수", BLUE),
        ("B", "감경", "금소법 과징금\n최대 75% 감경",
         "37.5억/건", "실태평가 우수 = Advocate 도입 근거", NAVY),
        ("C", "라이센스", "Phase 5 글로벌\n5개국 53개 은행",
         "1,000~3,000억/년", "$1M/은행 + customer fee", RED),
    ]
    card_w = Inches(4.0); card_h = Inches(4.4); gap = Inches(0.2)
    for i, (mark, label, what, big, sub, accent) in enumerate(axes):
        left = Inches(0.5) + (card_w + gap) * i
        top = Inches(2.1)
        add_card(s, left, top, card_w, card_h, fill=WHITE, border=accent, border_w=2.5)
        # 큰 알파벳 마커
        add_text(s, left + Inches(0.3), top + Inches(0.2), Inches(1.0), Inches(0.8),
                 text=f"({mark})", font_size=20, bold=True, color=accent,
                 font_name=H_FONT)
        add_text(s, left + Inches(0.3), top + Inches(0.85), Inches(3.4), Inches(0.6),
                 text=label, font_size=22, bold=True, color=NAVY,
                 font_name=H_FONT)
        add_text(s, left + Inches(0.3), top + Inches(1.5), Inches(3.4), Inches(0.9),
                 text=what, font_size=12, color=MUTE, font_name=H_FONT)
        add_text(s, left + Inches(0.3), top + Inches(2.55), Inches(3.4), Inches(0.8),
                 text=big, font_size=24, bold=True, color=accent,
                 font_name=H_FONT)
        add_text(s, left + Inches(0.3), top + Inches(3.4), Inches(3.4), Inches(0.9),
                 text=sub, font_size=11, color=INK, font_name=H_FONT)

    # 하단 break-even
    be = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5))
    be.fill.solid(); be.fill.fore_color.rgb = NAVY
    be.line.fill.background()
    add_text(s, Inches(0.7), Inches(6.78), Inches(11.9), Inches(0.4),
             text="누적 470억 ↔ Phase 5 라이센스 1.6년치로 회수",
             font_size=14, color=WHITE, bold=True, font_name=H_FONT,
             align=PP_ALIGN.CENTER)

    add_footer(s, 12)
    add_speaker_notes(s,
        "ROI 입니다. 3축 수익 모델. 첫째 회피 — "
        "DLF 197억급 사태를 1번만 회피해도 PoC 6개월 비용을 회수합니다. "
        "둘째 감경 — 금소법 과징금 최대 75% 감경 조건이 실태평가 우수입니다. "
        "Advocate 도입이 그 근거가 됩니다. "
        "셋째 라이센스 — Phase 5 글로벌 진출 시 싱가포르 MAS, 영국 FCA Consumer Duty 매핑으로 "
        "5개국 53개 은행 사용. 연 잠재 1,000억 ~ 3,000억.\n\n[45초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S13 — Paradigm Shift
# ──────────────────────────────────────────────────────────────────────────
def slide_13():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="패러다임 시프트", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="신기능이 아니라 역전입니다",
             font_size=30, bold=True, color=NAVY, font_name=H_FONT)

    shifts = [
        ("시점", "사후 모니터링", "사전 개입", "해피콜·녹취 분석 → 가입 직전 차단"),
        ("출력", "단일 추천", "이중 관점 동시", "Seller view + Advocate view 좌우 분할"),
        ("관점", "은행 편 AI", "고객 편 AI", "셀링 보조 → 고객 변호인"),
    ]
    for i, (axis, before, after, sub) in enumerate(shifts):
        top = Inches(2.3) + Inches(1.45) * i
        # 좌 (before)
        add_text(s, Inches(0.5), top, Inches(1.2), Inches(0.5),
                 text=axis, font_size=11, bold=True, color=MUTE, font_name=H_FONT)
        before_card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(1.8), top - Inches(0.05),
                                         Inches(3.7), Inches(1.0))
        before_card.fill.solid(); before_card.fill.fore_color.rgb = SOFT
        before_card.line.color.rgb = BORDER
        add_text(s, Inches(1.95), top + Inches(0.15), Inches(3.5), Inches(0.7),
                 text=before, font_size=18, color=MUTE, font_name=H_FONT,
                 align=PP_ALIGN.CENTER)

        # 화살표
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(5.7), top + Inches(0.2),
                                   Inches(1.4), Inches(0.6))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = RED
        arrow.line.fill.background()

        # 우 (after)
        after_card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(7.3), top - Inches(0.05),
                                        Inches(5.5), Inches(1.0))
        after_card.fill.solid(); after_card.fill.fore_color.rgb = NAVY
        after_card.line.fill.background()
        add_text(s, Inches(7.5), top + Inches(0.15), Inches(5.1), Inches(0.7),
                 text=after, font_size=20, bold=True, color=WHITE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        # sub line
        add_text(s, Inches(1.8), top + Inches(1.0), Inches(11), Inches(0.4),
                 text=sub, font_size=11, color=MUTE, font_name=H_FONT,
                 align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             text="\"고객 편 AI\" 라는 포지셔닝은 한 회사만 차지 — 우리은행",
             font_size=14, color=NAVY, bold=True,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 13)
    add_speaker_notes(s,
        "이건 단순한 신기능이 아니라 패러다임의 역전입니다. "
        "사후 → 사전, 단일 → 이중, 은행 → 고객. "
        "신한·KB·하나는 어쩌면 따라할 수 있을 것입니다. "
        "그러나 '고객 편 AI' 라는 포지셔닝은 한 회사만 차지합니다. "
        "그 한 회사가 우리은행입니다.\n\n[45초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S14 — ESG + Phase 4~5 글로벌
# ──────────────────────────────────────────────────────────────────────────
def slide_14():
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    add_text(s, Inches(0.5), Inches(0.4), Inches(8), Inches(0.6),
             text="Phase 4 ~ 5 · 글로벌", font_size=14, color=BLUE, bold=True, font_name=H_FONT)
    add_text(s, Inches(0.5), Inches(0.8), Inches(12), Inches(1.0),
             text="OCP Standard — 한국이 글로벌 reference 국가",
             font_size=28, bold=True, color=NAVY, font_name=H_FONT)

    # 상단 6단계 타임라인
    phases = [("P0", "설계", "3M"), ("P1", "PoC", "6M"), ("P2", "확장", "12M"),
              ("P3", "전사", "12M"), ("P4", "OCP 표준화", "24M"), ("P5", "글로벌", "24M")]
    box_w = Inches(2.0); gap = Inches(0.08); start_left = Inches(0.5)
    for i, (p, label, dur) in enumerate(phases):
        left = start_left + (box_w + gap) * i
        # accent: P4·P5 는 RED, 나머지 NAVY
        accent = RED if p in ("P4", "P5") else NAVY
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.1),
                               box_w, Inches(1.0))
        b.fill.solid(); b.fill.fore_color.rgb = accent
        b.line.fill.background()
        add_text(s, left, Inches(2.15), box_w, Inches(0.4),
                 text=p, font_size=14, bold=True, color=WHITE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, left, Inches(2.55), box_w, Inches(0.4),
                 text=label, font_size=12, color=WHITE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, left, Inches(2.95), box_w, Inches(0.3),
                 text=dur, font_size=10, color=LIGHT_BLUE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)

    # OCP 중앙 라벨
    add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.6),
             text="OCP — Open Consumer Protection Standard",
             font_size=22, bold=True, color=RED, font_name=H_FONT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.4),
             text="Linux Foundation 위탁 · 우리은행 최초 reference impl · 컨소시엄 의장사",
             font_size=12, color=MUTE, font_name=H_FONT, align=PP_ALIGN.CENTER)

    # 5개국 카드
    countries = [
        ("🇸🇬 싱가포르", "MAS FEAT Principles"),
        ("🇭🇰 홍콩", "SFC Suitability"),
        ("🇯🇵 일본", "금융청 · 시니어 보호"),
        ("🇪🇺 EU", "AI Act 첫 인증"),
        ("🇺🇸 미국", "SEC Reg BI"),
    ]
    cw = Inches(2.42); gap = Inches(0.05)
    for i, (flag, label) in enumerate(countries):
        left = Inches(0.5) + (cw + gap) * i
        add_card(s, left, Inches(4.8), cw, Inches(1.4), fill=SOFT, border=BORDER)
        add_text(s, left, Inches(4.95), cw, Inches(0.5),
                 text=flag, font_size=14, bold=True, color=NAVY,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)
        add_text(s, left, Inches(5.55), cw, Inches(0.7),
                 text=label, font_size=10, color=MUTE,
                 font_name=H_FONT, align=PP_ALIGN.CENTER)

    # ESG tag
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             text="ESG · S(취약계층 보호) + G(금소위 신설·표준 제정) → 우리은행 ESG 1등 추진",
             font_size=12, color=NAVY, bold=True,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_footer(s, 14)
    add_speaker_notes(s,
        "Phase 4~5 비전입니다. OCP — Open Consumer Protection Standard. "
        "우리은행이 한국 금융권 공통 소비자보호 AI 표준을 제정하고 Linux Foundation 에 위탁합니다. "
        "KB·신한이 자체 Advocate 를 구현해도 OCP 호환 인증을 우리은행에서 발급합니다. "
        "Phase 5 — 싱가포르 MAS FEAT, 영국 FCA Consumer Duty, EU AI Act 모두 OCP audit 으로 매핑됩니다. "
        "우리은행이 한국을 글로벌 금융 소비자보호 AI 의 reference 국가로 만듭니다.\n\n[45초]"
    )


# ──────────────────────────────────────────────────────────────────────────
# S15 — Conclusion (dark navy)
# ──────────────────────────────────────────────────────────────────────────
def slide_15():
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)

    add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.5),
             text="결언", font_size=14, color=LIGHT_BLUE, bold=True,
             font_name=H_FONT)

    # 3 lines, centered, large
    add_text(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.0),
             text="DLF 197억의 회사가",
             font_size=44, bold=True, color=WHITE,
             font_name=H_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.95), Inches(12.3), Inches(1.0),
             text="한국 금융권 소비자보호 표준의",
             font_size=44, bold=True, color=WHITE,
             font_name=H_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.0),
             text="출처가 됩니다.",
             font_size=54, bold=True, color=RED,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    # 하단
    sub = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.0))
    sub.fill.solid(); sub.fill.fore_color.rgb = RGBColor(0x00, 0x2A, 0x4B)
    sub.line.fill.background()
    add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.5),
             text="고객의 편에서 질문을 던지는",
             font_size=18, color=LIGHT_BLUE,
             font_name=H_FONT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.5),
             text="우리은행의 Woori Advocate 였습니다.",
             font_size=20, bold=True, color=WHITE,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
             text="Q&A — 들어주셔서 감사합니다",
             font_size=12, color=LIGHT_BLUE,
             font_name=H_FONT, align=PP_ALIGN.CENTER)

    add_speaker_notes(s,
        "마지막입니다.\n\n"
        "[★★ 3초 대기 후 천천히 한 줄씩 읽기 ★★]\n\n"
        "DLF 197억의 회사가\n"
        "한국 금융권 소비자보호 표준의\n"
        "출처가 됩니다.\n\n"
        "고객의 편에서 질문을 던지는 우리은행의 Woori Advocate 였습니다. "
        "들어주셔서 감사합니다.\n\n[45초] 마지막 한 호흡 쉰 뒤 인사."
    )


# ── Build all ────────────────────────────────────────────────────────────

builders = [
    slide_01, slide_02, slide_03, slide_04, slide_05,
    slide_06_mad_vs_langgraph,   # S6 신규 — MAD vs LangGraph
    slide_06,                     # S7 (구 S6, 학술 근거)
    slide_07,                     # S8 (구 S7, 규제)
    slide_08, slide_09, slide_10, slide_11, slide_12,
    slide_13, slide_14, slide_15,
]
# 페이지 번호 자동 갱신 (S6 추가로 총 16장)
# add_footer 의 total=15 기본값을 16으로 일괄 재호출 — 빌더 내부에서 idx 만 사용
for fn in builders:
    fn()

prs.save(OUT)
print(f"Built {len(prs.slides)} slides → {OUT}")
